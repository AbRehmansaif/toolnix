import os
import re
from django.conf import settings
from django.http import HttpResponse, JsonResponse
from django.views.decorators.csrf import csrf_exempt
from rest_framework.decorators import api_view
from django.core.files.storage import FileSystemStorage

from pdf2docx import Converter
from docx2pdf import convert as convert_docx_to_pdf

# PDF to Excel
import pdfplumber
import openpyxl

# Excel to PDF (Basic)
from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import letter

# PowerPoint
from pptx import Presentation

import tempfile
import uuid


def handle_upload(request):
    if request.method != 'POST' or not request.FILES.get('file'):
        return None, "No file uploaded"

    file = request.FILES['file']
    fs = FileSystemStorage(location=settings.MEDIA_ROOT)
    ext = os.path.splitext(file.name)[1].lower()
    filename = fs.save(f"{uuid.uuid4()}{ext}", file)
    return os.path.join(settings.MEDIA_ROOT, filename), file.name


def _pdf_to_docx_with_pymupdf(pdf_path, docx_path):
    """
    High-quality PDF → DOCX using PyMuPDF for text/font extraction +
    python-docx for document building.

    Preserves: font sizes, bold, italic, monospace, text colors,
    heading detection (by font size), paragraph spacing, page breaks.
    """
    import fitz  # PyMuPDF
    from docx import Document
    from docx.shared import Pt, Inches, RGBColor
    from docx.oxml.ns import qn
    from docx.oxml import OxmlElement

    doc = Document()

    # ── Standard margins ────────────────────────────────────────────────────
    for section in doc.sections:
        section.top_margin    = Inches(1.0)
        section.bottom_margin = Inches(1.0)
        section.left_margin   = Inches(1.25)
        section.right_margin  = Inches(1.25)

    # Remove Word's default empty paragraph
    for para in list(doc.paragraphs):
        para._element.getparent().remove(para._element)

    pdf_doc   = fitz.open(pdf_path)
    first_page = True

    for page in pdf_doc:
        # ── Page break between pages ────────────────────────────────────────
        if not first_page:
            p_elem = doc.add_paragraph()
            run    = p_elem.add_run()
            br     = OxmlElement('w:br')
            br.set(qn('w:type'), 'page')
            run._r.append(br)
        first_page = False

        blocks  = page.get_text("dict", sort=True).get("blocks", [])
        prev_y1 = None

        for block in blocks:
            if block.get("type") != 0:
                continue  # skip image blocks

            for line in block.get("lines", []):
                spans = line.get("spans", [])
                if not spans:
                    continue

                line_text = "".join(s.get("text", "") for s in spans).strip()
                if not line_text:
                    doc.add_paragraph()
                    prev_y1 = None
                    continue

                max_size = max((s.get("size", 11) for s in spans), default=11)

                # ── Add a blank line when there is a large vertical gap ─────
                y0 = line.get("bbox", [0, 0, 0, 0])[1]
                if prev_y1 is not None and (y0 - prev_y1) > max_size * 1.8:
                    doc.add_paragraph()
                prev_y1 = line.get("bbox", [0, 0, 0, 0])[3]

                para = doc.add_paragraph()

                # ── Heading detection by font size ─────────────────────────
                if max_size >= 20:
                    try:
                        para.style = doc.styles['Heading 1']
                    except Exception:
                        pass
                elif max_size >= 16:
                    try:
                        para.style = doc.styles['Heading 2']
                    except Exception:
                        pass
                elif max_size >= 13:
                    try:
                        para.style = doc.styles['Heading 3']
                    except Exception:
                        pass

                para.paragraph_format.space_after  = Pt(4)
                para.paragraph_format.space_before = Pt(2)

                for span in spans:
                    text = span.get("text", "")
                    if not text:
                        continue

                    font_size = span.get("size", 11)
                    flags     = span.get("flags", 0)
                    color_int = span.get("color", 0)
                    font_name = span.get("font", "")

                    run = para.add_run(text)
                    run.font.size   = Pt(round(font_size))
                    run.font.bold   = bool(flags & 16) or "Bold"   in font_name
                    run.font.italic = bool(flags & 2)  or "Italic" in font_name

                    # Monospace fonts
                    if (flags & 8) or any(
                        k in font_name for k in ("Mono", "Courier", "Code", "Fixed")
                    ):
                        run.font.name = "Courier New"

                    # Text color (skip pure black = 0)
                    if color_int and color_int != 0:
                        r = (color_int >> 16) & 0xFF
                        g = (color_int >> 8)  & 0xFF
                        b =  color_int        & 0xFF
                        if not (r == 0 and g == 0 and b == 0):
                            run.font.color.rgb = RGBColor(r, g, b)

    pdf_doc.close()
    doc.save(docx_path)


@csrf_exempt
@api_view(['POST'])
def pdf_to_word(request):
    input_path, original_name = handle_upload(request)
    if not input_path:
        return JsonResponse({'error': 'No file uploaded'}, status=400)

    base        = os.path.splitext(input_path)[0]
    output_path = base + '.docx'
    fallback    = base + '_v2.docx'

    safe_name = re.sub(r'[^\w\s\-.]', '', os.path.splitext(original_name)[0]).strip() or 'converted'

    def _send_docx(path):
        with open(path, 'rb') as f:
            resp = HttpResponse(
                f.read(),
                content_type='application/vnd.openxmlformats-officedocument.wordprocessingml.document',
            )
            resp['Content-Disposition'] = f'attachment; filename="{safe_name}.docx"'
            return resp

    try:
        # ── Stage 1: pdf2docx (best for tables & complex layouts) ──────────
        primary_ok = False
        try:
            cv = Converter(input_path)
            cv.convert(output_path, start=0, end=None)
            cv.close()
            # Consider it successful only if output is non-trivial
            if os.path.exists(output_path) and os.path.getsize(output_path) > 5_000:
                primary_ok = True
        except Exception:
            pass

        if primary_ok:
            return _send_docx(output_path)

        # ── Stage 2: PyMuPDF + python-docx (better text/font extraction) ───
        _pdf_to_docx_with_pymupdf(input_path, fallback)
        if os.path.exists(fallback) and os.path.getsize(fallback) > 1_000:
            return _send_docx(fallback)

        return JsonResponse(
            {'error': 'Could not convert this PDF. It may be scanned (image-only) or password-protected.'},
            status=422,
        )

    except Exception as exc:
        return JsonResponse({'error': f'Conversion failed: {exc}'}, status=500)

    finally:
        for p in (input_path, output_path, fallback):
            try:
                if os.path.exists(p):
                    os.remove(p)
            except Exception:
                pass



@csrf_exempt
@api_view(['POST'])
def word_to_pdf(request):
    input_path, original_name = handle_upload(request)
    if not input_path:
        return JsonResponse({'error': 'No file uploaded'}, status=400)
    
    output_path = input_path.replace('.docx', '.pdf').replace('.doc', '.pdf')
    try:
        # Requires MS Word installed on Windows
        convert_docx_to_pdf(input_path, output_path)
        
        with open(output_path, 'rb') as f:
            response = HttpResponse(f.read(), content_type='application/pdf')
            response['Content-Disposition'] = f'attachment; filename="{os.path.splitext(original_name)[0]}.pdf"'
            return response
    except Exception as e:
        return JsonResponse({'error': str(e)}, status=500)
    finally:
        if os.path.exists(input_path): os.remove(input_path)
        if os.path.exists(output_path): os.remove(output_path)

@csrf_exempt
@api_view(['POST'])
def excel_to_pdf(request):
    input_path, original_name = handle_upload(request)
    if not input_path:
        return JsonResponse({'error': 'No file uploaded'}, status=400)
    
    output_path = input_path.replace('.xlsx', '.pdf').replace('.xls', '.pdf')
    try:
        wb = openpyxl.load_workbook(input_path, data_only=True)
        sheet = wb.active
        
        c = canvas.Canvas(output_path, pagesize=letter)
        y = 750
        for row in sheet.iter_rows(values_only=True):
            text = " | ".join([str(cell) if cell is not None else "" for cell in row])
            c.drawString(50, y, text[:100]) # simple text render
            y -= 20
            if y < 50:
                c.showPage()
                y = 750
        c.save()
        
        with open(output_path, 'rb') as f:
            response = HttpResponse(f.read(), content_type='application/pdf')
            response['Content-Disposition'] = f'attachment; filename="{os.path.splitext(original_name)[0]}.pdf"'
            return response
    except Exception as e:
        return JsonResponse({'error': str(e)}, status=500)
    finally:
        if os.path.exists(input_path): os.remove(input_path)
        if os.path.exists(output_path): os.remove(output_path)

@csrf_exempt
@api_view(['POST'])
def pdf_to_excel(request):
    input_path, original_name = handle_upload(request)
    if not input_path:
        return JsonResponse({'error': 'No file uploaded'}, status=400)
    
    output_path = input_path.replace('.pdf', '.xlsx')
    try:
        wb = openpyxl.Workbook()
        ws = wb.active
        
        with pdfplumber.open(input_path) as pdf:
            for page in pdf.pages:
                text = page.extract_text()
                if text:
                    for line in text.split('\n'):
                        ws.append([line])
                        
        wb.save(output_path)
        
        with open(output_path, 'rb') as f:
            response = HttpResponse(f.read(), content_type='application/vnd.openxmlformats-officedocument.spreadsheetml.sheet')
            response['Content-Disposition'] = f'attachment; filename="{os.path.splitext(original_name)[0]}.xlsx"'
            return response
    except Exception as e:
        return JsonResponse({'error': str(e)}, status=500)
    finally:
        if os.path.exists(input_path): os.remove(input_path)
        if os.path.exists(output_path): os.remove(output_path)

@csrf_exempt
@api_view(['POST'])
def pptx_to_pdf(request):
    input_path, original_name = handle_upload(request)
    if not input_path:
        return JsonResponse({'error': 'No file uploaded'}, status=400)
    
    output_path = input_path.replace('.pptx', '.pdf').replace('.ppt', '.pdf')
    try:
        prs = Presentation(input_path)
        c = canvas.Canvas(output_path, pagesize=letter)
        for slide in prs.slides:
            y = 700
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    c.drawString(50, y, shape.text[:100])
                    y -= 20
            c.showPage()
        c.save()
        
        with open(output_path, 'rb') as f:
            response = HttpResponse(f.read(), content_type='application/pdf')
            response['Content-Disposition'] = f'attachment; filename="{os.path.splitext(original_name)[0]}.pdf"'
            return response
    except Exception as e:
        return JsonResponse({'error': str(e)}, status=500)
    finally:
        if os.path.exists(input_path): os.remove(input_path)
        if os.path.exists(output_path): os.remove(output_path)

@csrf_exempt
@api_view(['POST'])
def pdf_to_pptx(request):
    input_path, original_name = handle_upload(request)
    if not input_path:
        return JsonResponse({'error': 'No file uploaded'}, status=400)
    
    output_path = input_path.replace('.pdf', '.pptx')
    try:
        prs = Presentation()
        title_slide_layout = prs.slide_layouts[0]
        
        with pdfplumber.open(input_path) as pdf:
            for page in pdf.pages:
                slide = prs.slides.add_slide(title_slide_layout)
                title = slide.shapes.title
                title.text = page.extract_text()[:50] if page.extract_text() else "Slide"
                
        prs.save(output_path)
        
        with open(output_path, 'rb') as f:
            response = HttpResponse(f.read(), content_type='application/vnd.openxmlformats-officedocument.presentationml.presentation')
            response['Content-Disposition'] = f'attachment; filename="{os.path.splitext(original_name)[0]}.pptx"'
            return response
    except Exception as e:
        return JsonResponse({'error': str(e)}, status=500)
    finally:
        if os.path.exists(input_path): os.remove(input_path)
        if os.path.exists(output_path): os.remove(output_path)
